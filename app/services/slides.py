import asyncio
import io
import json
import math
import re
import uuid

from pydantic import ValidationError

from app.schemas.language import OutputLanguage
from app.schemas.slides import SlideDeckResponse, SlideItem
from app.services.llm.factory import get_llm_provider
from app.services.llm.types import Message
from app.services.pexels_stock import fetch_landscape_photo_bytes

_lock = asyncio.Lock()
_decks: dict[str, tuple[bytes, str]] = {}


def _lang_rule(lang: OutputLanguage) -> str:
    if lang == "english":
        return "All slide text must be in English."
    if lang == "hindi":
        return "All slide text must be in Hindi (Devanagari script)."
    return "All slide text must be in Roman Hindi (Latin letters only, no Devanagari)."


def _prompt(
    *,
    topic: str,
    slide_count: int,
    output_language: OutputLanguage,
    context_text: str | None,
) -> str:
    ctx = (context_text or "").strip()
    ctx_excerpt = ctx[:26000]
    return (
        "You are an expert teacher preparing a classroom slide deck.\n"
        f"Topic: {topic}\n"
        f"Number of content slides (not counting the opening): {slide_count}\n"
        f"{_lang_rule(output_language)}\n\n"
        "Return JSON only (no markdown). Keys:\n"
        '- deck_title: short presentation title (max 80 characters)\n'
        f"- slides: array of exactly {slide_count} objects. Each object has:\n"
        "  title (slide heading), bullets (array of 3 to 6 concise strings), optional speaker_notes (teacher script),\n"
        "  and image_query (English only, REQUIRED on every slide): a vivid, SPECIFIC stock-photo search phrase "
        "(4–8 words) derived from the actual slide title and bullet content. "
        "The query must be a concrete visual scene — a real-world subject, setting, or action — "
        "NOT abstract concepts. Examples: "
        "\"astronaut floating in space station\", "
        "\"scientist examining DNA strands in lab\", "
        "\"ancient roman aqueduct stone arch\", "
        "\"child reading book under tree sunlight\". "
        "Bad examples (too generic, do NOT use): \"education classroom learning\", \"science concept background\", \"technology abstract\". "
        "Make the query unique per slide — reflect what THAT slide is specifically about.\n\n"
        "Rules:\n"
        "- Bullets must be teaching points, not slide design instructions.\n"
        "- Order slides as you would teach: intro → concepts → examples → summary.\n"
        "- Ground content in the Context when it is provided.\n"
        "- No HTML or markdown in strings.\n"
        + (f"\nContext material:\n{ctx_excerpt}\n" if ctx_excerpt else "")
    )


def _extract_json_obj(raw: str) -> dict | None:
    raw = raw.strip()
    try:
        obj = json.loads(raw)
        return obj if isinstance(obj, dict) else None
    except Exception:
        pass
    start = raw.find("{")
    end = raw.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            obj = json.loads(raw[start : end + 1])
            return obj if isinstance(obj, dict) else None
        except Exception:
            return None
    return None


def _normalize_slides(raw_slides: list, *, expected: int, topic: str) -> list[SlideItem]:
    out: list[SlideItem] = []
    for i, item in enumerate(raw_slides[:expected]):
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", f"Slide {i + 1}")).strip()[:200]
        bullets_raw = item.get("bullets") or []
        if not isinstance(bullets_raw, list):
            bullets_raw = [str(bullets_raw)]
        bullets = [str(b).strip()[:500] for b in bullets_raw if str(b).strip()][:8]
        if not bullets:
            bullets = [f"Key idea {i + 1} about {topic}."]
        notes = item.get("speaker_notes")
        sn = str(notes).strip()[:4000] if notes else None
        iq = item.get("image_query")
        image_query = str(iq).strip()[:200] if iq else None
        try:
            out.append(
                SlideItem(
                    title=title or f"Slide {i + 1}",
                    bullets=bullets,
                    speaker_notes=sn,
                    image_query=image_query or None,
                )
            )
        except ValidationError:
            continue
    while len(out) < expected:
        n = len(out)
        out.append(
            SlideItem(
                title=f"{topic} — point {n + 1}",
                bullets=[f"Review this section in the chapter.", "Ask students for one example."],
                speaker_notes=None,
                image_query=None,
            )
        )
    return out[:expected]


def _synthetic_slide_background(slide_index: int, total: int) -> bytes:
    """16:9 PNG — cinematic gradient + accents."""
    from PIL import Image, ImageDraw

    w, h = 1600, 900
    if slide_index < 0:
        top, bot = (12, 18, 42), (28, 64, 112)
        accent = (56, 189, 248)
    else:
        palettes = (
            ((16, 24, 40), (36, 52, 88), (99, 102, 241)),
            ((18, 32, 48), (42, 78, 96), (45, 212, 191)),
            ((26, 20, 46), (62, 40, 86), (192, 132, 252)),
            ((20, 28, 36), (48, 58, 72), (251, 191, 36)),
        )
        top, bot, accent = palettes[slide_index % len(palettes)]

    img = Image.new("RGB", (w, h), top)
    draw = ImageDraw.Draw(img)
    for y in range(h):
        t = y / max(h - 1, 1)
        r = int(top[0] + (bot[0] - top[0]) * t)
        g = int(top[1] + (bot[1] - top[1]) * t)
        b = int(top[2] + (bot[2] - top[2]) * t)
        draw.line([(0, y), (w, y)], fill=(r, g, b))

    cx = int(w * (0.82 + 0.06 * math.sin(slide_index * 0.9)))
    cy = int(h * (0.28 + 0.12 * math.cos(slide_index * 0.7)))
    for rad, width in ((260, 5), (420, 3)):
        bbox = (cx - rad, cy - rad, cx + rad, cy + rad)
        draw.ellipse(bbox, outline=accent, width=width)

    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def _prepare_stock_photo_for_slide(raw: bytes) -> bytes:
    """Resize to 16:9 right-panel size and blend a dark veil."""
    from PIL import Image

    im = Image.open(io.BytesIO(raw)).convert("RGB")
    im = im.resize((800, 900), Image.Resampling.LANCZOS)  # right-half panel size
    veil = Image.new("RGB", im.size, (8, 11, 22))
    blended = Image.blend(im, veil, 0.25)  # lighter veil so image is more visible
    buf = io.BytesIO()
    blended.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def _build_pptx_bytes(
    *,
    deck_title: str,
    topic: str,
    slides: list[SlideItem],
    slide_images: list[bytes],
) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    # Fix: set correct 16:9 widescreen dimensions
    prs.slide_width = Emu(12192000)   # 13.33 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    sw, sh = prs.slide_width, prs.slide_height

    try:
        prs.core_properties.title = deck_title[:200]
    except Exception:
        pass

    blank_idx = 6
    try:
        blank_layout = prs.slide_layouts[blank_idx]
    except Exception:
        blank_layout = prs.slide_layouts[1]

    def _full_bleed_picture(slide, png_bytes: bytes) -> None:
        stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(stream, 0, 0, width=sw, height=sh)

    def _right_panel_picture(slide, png_bytes: bytes) -> None:
        """Place stock photo only on the right half of the slide."""
        stream = io.BytesIO(png_bytes)
        left = Inches(5.5)
        slide.shapes.add_picture(stream, left, 0, width=sw - left, height=sh)

    def _rgb(r: int, g: int, b: int) -> RGBColor:
        return RGBColor(r, g, b)

    # --- Opening "title" slide ---
    hero = _synthetic_slide_background(-1, max(len(slides), 1))
    s0 = prs.slides.add_slide(blank_layout)
    _full_bleed_picture(s0, hero)
    veil0 = s0.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    veil0.fill.solid()
    veil0.fill.fore_color.rgb = _rgb(6, 8, 18)
    veil0.fill.transparency = 0.42
    veil0.line.fill.background()

    accent_bar = s0.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(2.15), Inches(0.12), Inches(2.2))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _rgb(56, 189, 248)
    accent_bar.line.fill.background()

    tb0 = s0.shapes.add_textbox(Inches(1.1), Inches(2.05), Inches(11.2), Inches(1.35))
    tf0 = tb0.text_frame
    tf0.word_wrap = True
    tf0.text = deck_title[:200]
    p0 = tf0.paragraphs[0]
    p0.font.size = Pt(40)
    p0.font.bold = True
    p0.font.color.rgb = _rgb(248, 250, 252)
    p0.line_spacing = 1.05

    sub0 = s0.shapes.add_textbox(Inches(1.1), Inches(3.55), Inches(11.0), Inches(1.1))
    sf0 = sub0.text_frame
    sf0.text = topic[:320]
    sp0 = sf0.paragraphs[0]
    sp0.font.size = Pt(18)
    sp0.font.color.rgb = _rgb(186, 198, 216)

    foot = s0.shapes.add_textbox(Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.55))
    ff = foot.text_frame
    ff.text = "ClassroomAI"
    fp = ff.paragraphs[0]
    fp.font.size = Pt(11)
    fp.font.color.rgb = _rgb(148, 163, 184)

    # --- Content slides ---
    for idx, item in enumerate(slides):
        img_b = slide_images[idx] if idx < len(slide_images) else _synthetic_slide_background(idx, len(slides))
        slide = prs.slides.add_slide(blank_layout)

        has_stock = bool(item.photo_attribution)

        if has_stock:
            # Dark background base
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
            bg.fill.solid()
            bg.fill.fore_color.rgb = _rgb(11, 15, 28)
            bg.line.fill.background()

            # Stock photo on right half only — now actually visible
            _right_panel_picture(slide, img_b)

            # Dark panel over left half for text readability
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.5), sh)
            panel.fill.solid()
            panel.fill.fore_color.rgb = _rgb(11, 15, 28)
            panel.fill.transparency = 0.08
            panel.line.fill.background()

            title_left = Inches(0.45)
            title_w = Inches(4.85)
            body_left = Inches(0.45)
            body_w = Inches(4.85)
            body_top = Inches(1.45)
            body_h = Inches(5.2)
        else:
            # Synthetic gradient full bleed
            _full_bleed_picture(slide, img_b)
            leg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.3), sw, Inches(3.2))
            leg.fill.solid()
            leg.fill.fore_color.rgb = _rgb(8, 10, 20)
            leg.fill.transparency = 0.28
            leg.line.fill.background()
            title_left = Inches(0.55)
            title_w = Inches(12.2)
            body_left = Inches(0.65)
            body_w = Inches(12.0)
            body_top = Inches(1.35)
            body_h = Inches(5.2)

        # Title
        tbox = slide.shapes.add_textbox(title_left, Inches(0.32), title_w, Inches(0.95))
        ttf = tbox.text_frame
        ttf.text = item.title[:250]
        tp = ttf.paragraphs[0]
        tp.font.size = Pt(28) if has_stock else Pt(30)
        tp.font.bold = True
        tp.font.color.rgb = _rgb(248, 250, 252)

        # Bullets
        bbox = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
        tf = bbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        bullets = item.bullets[:8] or ["—"]
        tf.text = bullets[0][:500]
        tf.paragraphs[0].font.size = Pt(16 if has_stock else 17)
        tf.paragraphs[0].font.color.rgb = _rgb(214, 222, 235)
        tf.paragraphs[0].space_after = Pt(8)
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b[:500]
            p.level = 0
            p.font.size = Pt(15 if has_stock else 16)
            p.font.color.rgb = _rgb(196, 208, 224)
            p.space_after = Pt(6)

        # Speaker notes
        notes_parts: list[str] = []
        if item.speaker_notes:
            notes_parts.append(item.speaker_notes[:4000])
        if item.photo_attribution:
            notes_parts.append(item.photo_attribution)
        note_text = "\n\n".join(notes_parts).strip()
        if note_text:
            try:
                ns = slide.notes_slide.notes_text_frame
                ns.text = note_text[:4500]
            except Exception:
                pass

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _safe_filename(name: str) -> str:
    s = re.sub(r'[<>:"/\\|?*]', "", name).strip()
    s = re.sub(r"\s+", " ", s)[:100]
    return s or "ClassroomAI-slides"


async def generate_slide_deck(
    *,
    topic: str,
    context_text: str | None,
    output_language: OutputLanguage,
    slide_count: int,
    llm_provider: str,
) -> SlideDeckResponse:
    from app.core.config import get_settings

    provider = get_llm_provider(vision=False, provider=llm_provider)
    messages: list[Message] = [
        {"role": "system", "content": "You output strict JSON only for slide deck generation."},
        {
            "role": "user",
            "content": _prompt(
                topic=topic,
                slide_count=slide_count,
                output_language=output_language,
                context_text=context_text,
            ),
        },
    ]
    raw = await provider.complete_chat(messages, max_tokens=8000, temperature=0.35)
    data = _extract_json_obj(raw) or {}
    deck_title = str(data.get("deck_title", topic)).strip()[:200] or topic
    raw_slides = data.get("slides")
    if not isinstance(raw_slides, list):
        raw_slides = []
    slides = _normalize_slides(raw_slides, expected=slide_count, topic=topic)

    settings = get_settings()
    pexels_key = (settings.pexels_api_key or "").strip()
    attrs: list[str | None] = [None] * len(slides)
    slide_pngs: list[bytes] = []
    for i, s in enumerate(slides):
        raw: bytes | None = None
        if pexels_key:
            q = (s.image_query or topic).strip()
            out = await fetch_landscape_photo_bytes(pexels_key, q)
            # Fallback 1: retry with just the slide title if image_query fails
            if not out:
                fallback_q = s.title.strip()
                out = await fetch_landscape_photo_bytes(pexels_key, fallback_q)
            # Fallback 2: use the topic itself
            if not out:
                out = await fetch_landscape_photo_bytes(pexels_key, topic)
            if out:
                raw, _mime, attr = out
                attrs[i] = attr
        if raw:
            slide_pngs.append(_prepare_stock_photo_for_slide(raw))
        else:
            slide_pngs.append(_synthetic_slide_background(i, len(slides)))

    final_slides = [
        s.model_copy(update={"photo_attribution": attrs[i]})
        for i, s in enumerate(slides)
    ]

    pptx_bytes = _build_pptx_bytes(
        deck_title=deck_title,
        topic=topic,
        slides=final_slides,
        slide_images=slide_pngs,
    )
    deck_id = f"sdeck_{uuid.uuid4().hex[:16]}"
    filename = f"{_safe_filename(deck_title)}.pptx"

    async with _lock:
        _decks[deck_id] = (pptx_bytes, filename)

    return SlideDeckResponse(
        deck_id=deck_id,
        deck_title=deck_title,
        filename=filename,
        slides=final_slides,
    )


async def get_deck_file(deck_id: str) -> tuple[bytes, str] | None:
    async with _lock:
        return _decks.get(deck_id)